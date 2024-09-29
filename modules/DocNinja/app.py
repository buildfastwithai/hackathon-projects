import os
import requests
import streamlit as st
import fitz  
import docx
import pandas as pd
from bs4 import BeautifulSoup
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationChain
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from pptx import Presentation  

load_dotenv()
google_api_key = os.getenv("GOOGLE_API_KEY")

# Function to extract text from a PDF file
def extract_text_from_pdf(file):
    text = ""
    pdf_document = fitz.open(file)
    for page in pdf_document:
        text += page.get_text()
    pdf_document.close()
    return text

# Function to extract text from a DOCX file
def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text.strip()

# Function to extract text from a TXT file
def extract_text_from_txt(file):
    return file.read().decode("utf-8")

# Function to extract text from an HTML file
def extract_text_from_html(file):
    soup = BeautifulSoup(file.read(), "html.parser")
    return soup.get_text()

# Function to extract text from a CSV file
def extract_text_from_csv(file):
    df = pd.read_csv(file)
    return df.to_string(index=False)

# Function to extract text from a PPTX file
def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

# Function to summarize text and prepare for Q&A
def summarize_and_answer(text):
    llm = ChatGoogleGenerativeAI(
        model="gemini-1.5-flash",
        temperature=0,
        max_tokens=None,
        timeout=None,
        max_retries=2,
        api_key=google_api_key
    )

    # Summarization
    messages_summary = [
        ("system", "You're a helpful assistant that summarizes documents."),
        ("human", text),
    ]
    summary = llm.invoke(messages_summary).content

    # Prepare for Q&A
    memory = ConversationBufferMemory()
    conversation = ConversationChain(
        llm=llm,
        verbose=False,
        memory=memory
    )

    # Save initial context for Q&A
    memory.save_context({"input": text}, {"output": ""})

    return summary, conversation

# Function to get text from PDF
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text() or ""
    return text

# Function to split text into chunks for processing
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

# Function to create a vector store from text chunks
def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# Function to create a conversational chain for LLM Q&A
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. If the answer is not in provided context, just say, "The context is not provided in the document." Please do not provide a wrong answer.\n\n
    Context:\n {context}?\n
    Question: \n{question}\n
    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)

    return chain

# Main function for the Streamlit app
def main():
    # Set page config
    st.set_page_config(page_title="DocuNinja", layout="wide")

    # Add title and description
    st.title("🗂️ DocuNinja")
    st.write("Master your document review process with **DocuNinja**! Effortlessly upload any document format—be it **PPTs, DOCs, PDFs, Txts, HTMLs, or CSVs**—and let our advanced AI transform lengthy texts into concise summaries. With integrated Q&A capabilities, **DocuNinja** helps you clarify complex ideas and extract critical insights swiftly. Whether you're a student preparing for exams or a professional needing quick insights, **DocuNinja** equips you with the tools to maximize the value of your documents in record time!")

    # Add a sidebar for file upload and styling
    st.sidebar.header("Upload Document")
    uploaded_file = st.sidebar.file_uploader("Choose a PDF, DOCX, TXT, HTML, CSV, or PPTX file", type=["pdf", "docx", "txt", "html", "csv", "pptx"])

    if uploaded_file is not None:
        # Extract text based on file type
        if uploaded_file.type == "application/pdf":
            text = extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            text = extract_text_from_docx(uploaded_file)
        elif uploaded_file.type == "text/plain":
            text = extract_text_from_txt(uploaded_file)
        elif uploaded_file.type == "text/html":
            text = extract_text_from_html(uploaded_file)
        elif uploaded_file.type == "text/csv":
            text = extract_text_from_csv(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = extract_text_from_pptx(uploaded_file)

        if text:
            summary, conversation = summarize_and_answer(text)

            # Display the summary in a card format
            st.subheader("Document Summary:")
            st.markdown(f"<div style='border: 1px solid #e0e0e0; border-radius: 5px; padding: 15px; background-color: #f9f9f9;'>{summary}</div>", unsafe_allow_html=True)

            # Initialize session state for chat messages
            if "doc_messages" not in st.session_state:
                st.session_state.doc_messages = []

            # Document Q&A section
            st.subheader("💬 Ask Questions About the Document")
            doc_messages = st.container()
            with doc_messages:
                # Display chat messages
                for message in st.session_state.doc_messages:
                    if message['role'] == 'user':
                        st.chat_message("user").write(message['content'])
                    else:
                        st.chat_message("assistant").write(message['content'])

                # Get user input for document-related questions
                if prompt := st.chat_input("Ask a question about the document"):
                    st.session_state.doc_messages.append({"role": "user", "content": prompt})
                    st.chat_message("user").write(prompt)

                    with st.spinner("Generating response..."):
                        response = conversation.predict(input=prompt)
                        if response.strip() == "":
                            response = "The context is not provided in the document."

                    st.session_state.doc_messages.append({"role": "assistant", "content": response})
                    st.chat_message("assistant").write(response)


if __name__ == "__main__":
    main()
